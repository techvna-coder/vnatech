import streamlit as st
from io import BytesIO
import time
import re
from typing import List, Dict, Any, Tuple

# OpenAI (cho embeddings)
try:
    from openai import OpenAI
    client = OpenAI(api_key=st.secrets.get("OPENAI_API_KEY", ""))
except ImportError:
    st.error("OpenAI library not installed")
    client = None

# PDF/PPTX
try:
    import PyPDF2
except ImportError:
    st.error("PyPDF2 not installed")
    PyPDF2 = None

try:
    from pptx import Presentation
except ImportError:
    st.error("python-pptx not installed")
    Presentation = None

# Tokenizer
try:
    import tiktoken
    tokenizer = tiktoken.encoding_for_model("text-embedding-3-small")
except Exception:
    tokenizer = None
    st.warning("tiktoken not available; chunking will fall back to rough splitting.")

# ---------- Utilities ----------
def _safe_tokenize(text: str) -> List[int]:
    if tokenizer:
        return tokenizer.encode(text)
    # fallback: giả lập 'token' theo ký tự
    return list(text)

def _safe_detokenize(tokens: List[int]) -> str:
    if tokenizer:
        return tokenizer.decode(tokens)
    # fallback: ghép lại ký tự
    return "".join(tokens)

def count_tokens(text: str) -> int:
    try:
        return len(_safe_tokenize(text))
    except Exception:
        return 0

def preprocess_text(text: str) -> str:
    # Chuẩn hóa nhẹ, tránh phá vỡ nội dung kỹ thuật
    try:
        text = text.replace("\f", "\n").replace("\r", "\n")
        # Chuẩn hóa dấu ngoặc kép và gạch ngang thường gặp từ OCR
        text = text.replace("–", "-").replace("—", "-")
        text = text.replace(""", '"').replace(""", '"')
        text = text.replace("'", "'").replace("'", "'")
        # Dọn khoảng trắng thừa theo dòng, giữ xuống dòng hợp lý
        lines = [ln.strip() for ln in text.split("\n")]
        text = "\n".join([ln for ln in lines if ln])
        return text.strip()
    except Exception:
        return text

# ---------- Extractors with metadata ----------
def process_pdf(file_content: BytesIO) -> Tuple[str, Dict[str, Any]]:
    """Extract text + metadata from PDF."""
    if PyPDF2 is None:
        raise Exception("PyPDF2 not installed")
    try:
        reader = PyPDF2.PdfReader(file_content)
        pages = []
        for i, page in enumerate(reader.pages, 1):
            try:
                t = page.extract_text() or ""
            except Exception as e:
                st.warning(f"Failed to extract text from page {i}: {e}")
                t = ""
            if t.strip():
                pages.append(f"--- Page {i} ---\n{t}")
        full_text = "\n".join(pages).strip()
        if not full_text:
            raise Exception("No text could be extracted from the PDF")
        meta = {
            "total_pages": len(reader.pages),
            "has_sections": True,
            "sections": [{"type": "page", "number": i} for i in range(1, len(reader.pages)+1)],
            "has_tables": False,  # placeholder
            "has_lists": False,   # placeholder
            "key_terms": []
        }
        return full_text, meta
    except Exception as e:
        raise Exception(f"Failed to process PDF: {str(e)}")

def process_pptx(file_content: BytesIO) -> Tuple[str, Dict[str, Any]]:
    """Extract text + metadata from PPTX."""
    if Presentation is None:
        raise Exception("python-pptx not installed")
    try:
        prs = Presentation(file_content)
        slides_text = []
        has_tables = False
        
        for s_idx, slide in enumerate(prs.slides, 1):
            buf = [f"--- Slide {s_idx} ---"]
            for shape in slide.shapes:
                try:
                    # Textbox hoặc text frame
                    if hasattr(shape, "text") and shape.text:
                        buf.append(shape.text)
                    
                    # Table - kiểm tra kỹ càng hơn
                    if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE = 19
                        has_tables = True
                        if hasattr(shape, "table"):
                            table = shape.table
                            for row in table.rows:
                                cells = []
                                for cell in row.cells:
                                    try:
                                        cell_text = cell.text.strip()
                                        if cell_text:
                                            cells.append(cell_text)
                                    except Exception:
                                        continue
                                if cells:
                                    buf.append(" | ".join(cells))
                except Exception as e:
                    # Bỏ qua shape lỗi, tiếp tục xử lý
                    st.warning(f"Skipped shape in slide {s_idx}: {str(e)}")
                    continue
            
            slide_block = "\n".join([b for b in buf if b and b.strip() != f"--- Slide {s_idx} ---"])
            if slide_block.strip():
                slides_text.append(slide_block)
        
        full_text = "\n".join(slides_text).strip()
        if not full_text:
            raise Exception("No text could be extracted from the PowerPoint")
        
        meta = {
            "total_slides": len(list(prs.slides)),
            "has_sections": True,
            "sections": [{"type": "slide", "number": i} for i in range(1, len(list(prs.slides))+1)],
            "has_tables": has_tables,
            "has_lists": False,
            "key_terms": []
        }
        return full_text, meta
    except Exception as e:
        raise Exception(f"Failed to process PPTX: {str(e)}")

# ---------- Smart Chunking with Semantic Boundaries ----------
def _detect_natural_breaks(text: str) -> List[int]:
    """
    Phát hiện vị trí các ranh giới tự nhiên trong text:
    - Dòng trống (paragraph breaks)
    - Headers (dòng ngắn kết thúc bằng dấu hai chấm hoặc số)
    - Bullet points
    - Numbered lists
    
    Returns: danh sách các vị trí có thể cắt (character indices)
    """
    breaks = [0]  # Bắt đầu
    lines = text.split('\n')
    pos = 0
    
    for i, line in enumerate(lines):
        pos += len(line) + 1  # +1 cho \n
        
        # Dòng trống - ranh giới mạnh
        if not line.strip():
            breaks.append(pos)
            continue
        
        # Header patterns (dòng ngắn kết thúc bằng : hoặc số)
        if len(line.strip()) < 100:
            # Kết thúc bằng dấu hai chấm
            if line.strip().endswith(':'):
                breaks.append(pos)
                continue
            
            # Bắt đầu bằng số (numbered section)
            if re.match(r'^\s*\d+[\.\)]\s+', line):
                breaks.append(pos)
                continue
            
            # Header style: "ATA XX-XX -"
            if re.match(r'^\s*[A-Z0-9]+-?\s*[A-Z0-9]*\s*-', line):
                breaks.append(pos)
                continue
        
        # Bullet points
        if re.match(r'^\s*[-•*]\s+', line):
            breaks.append(pos)
            continue
    
    breaks.append(len(text))  # Kết thúc
    return sorted(set(breaks))

def _find_best_split_point(text: str, target_pos: int, window: int = 200) -> int:
    """
    Tìm điểm cắt tốt nhất gần target_pos trong khoảng window.
    Ưu tiên: paragraph break > sentence end > word boundary
    """
    start = max(0, target_pos - window)
    end = min(len(text), target_pos + window)
    search_zone = text[start:end]
    
    # 1. Tìm paragraph break (double newline hoặc single newline với dòng trống)
    para_breaks = [m.end() for m in re.finditer(r'\n\s*\n', search_zone)]
    if para_breaks:
        closest = min(para_breaks, key=lambda x: abs(x - (target_pos - start)))
        return start + closest
    
    # 2. Tìm sentence end (. ! ?) followed by space/newline
    sent_breaks = [m.end() for m in re.finditer(r'[.!?]\s+', search_zone)]
    if sent_breaks:
        closest = min(sent_breaks, key=lambda x: abs(x - (target_pos - start)))
        return start + closest
    
    # 3. Tìm word boundary
    word_breaks = [m.end() for m in re.finditer(r'\s+', search_zone)]
    if word_breaks:
        closest = min(word_breaks, key=lambda x: abs(x - (target_pos - start)))
        return start + closest
    
    # Fallback: cắt tại target_pos
    return target_pos

def _chunk_by_semantic_boundaries(text: str, chunk_size: int, chunk_overlap: int) -> List[str]:
    """
    Chunk text theo semantic boundaries thay vì cắt cứng theo token.
    """
    tokens = _safe_tokenize(text)
    
    if len(tokens) <= chunk_size:
        return [text]
    
    # Lấy các ranh giới tự nhiên
    natural_breaks = _detect_natural_breaks(text)
    
    chunks = []
    start_char = 0
    
    while start_char < len(text):
        # Tính vị trí kết thúc mong muốn (theo token)
        start_tokens = _safe_tokenize(text[:start_char])
        target_tokens = len(start_tokens) + chunk_size
        
        # Tìm vị trí character tương ứng (xấp xỉ)
        if target_tokens >= len(tokens):
            end_char = len(text)
        else:
            # Binary search để tìm vị trí char tương ứng với token count
            left, right = start_char, len(text)
            while left < right:
                mid = (left + right) // 2
                mid_tokens = len(_safe_tokenize(text[:mid]))
                if mid_tokens < target_tokens:
                    left = mid + 1
                else:
                    right = mid
            end_char = left
        
        # Tìm natural break gần nhất
        natural_candidates = [b for b in natural_breaks if start_char < b <= end_char + 100]
        if natural_candidates:
            # Chọn break gần end_char nhất
            end_char = min(natural_candidates, key=lambda x: abs(x - end_char))
        else:
            # Không có natural break, tìm sentence/word boundary
            end_char = _find_best_split_point(text, end_char, window=150)
        
        # Extract chunk
        chunk_text = text[start_char:end_char].strip()
        if chunk_text:
            chunks.append(chunk_text)
        
        if end_char >= len(text):
            break
        
        # Tính overlap theo token
        overlap_tokens = min(chunk_overlap, len(_safe_tokenize(chunk_text)) // 2)
        if overlap_tokens > 0:
            # Lùi lại overlap_tokens
            overlap_text = _safe_detokenize(_safe_tokenize(chunk_text)[-overlap_tokens:])
            # Tìm vị trí overlap trong text gốc
            overlap_pos = text.rfind(overlap_text[:50], start_char, end_char)
            if overlap_pos > start_char:
                start_char = overlap_pos
            else:
                start_char = end_char
        else:
            start_char = end_char
    
    return chunks

def chunk_text_smart(text: str,
                     doc_metadata: Dict[str, Any],
                     chunk_size: int = 1000,
                     chunk_overlap: int = 200) -> List[Dict[str, Any]]:
    """
    Tách theo 'đơn vị logic' (Page/Slide) nếu có, sau đó chunk theo semantic boundaries.
    """
    text = preprocess_text(text)
    sections = []
    
    # Ưu tiên cắt theo nhãn đã chèn trong extractor
    split_pages = [p for p in text.split("\n--- Page ") if p.strip()]
    split_slides = [s for s in text.split("\n--- Slide ") if s.strip()]

    if len(split_pages) > 1:  # PDF
        for blk in split_pages:
            try:
                header, *rest = blk.split("---", 1)
                number = int(header.strip())
                content = rest[0] if rest else ""
            except Exception:
                number, content = 0, blk
            sections.append(("page", number, content.strip()))
    elif len(split_slides) > 1:  # PPTX
        for blk in split_slides:
            try:
                header, *rest = blk.split("---", 1)
                number = int(header.strip())
                content = rest[0] if rest else ""
            except Exception:
                number, content = 0, blk
            sections.append(("slide", number, content.strip()))
    else:
        # fallback: cả tài liệu 1 khối
        sections.append(("document", 0, text))

    out: List[Dict[str, Any]] = []
    total_chunks_counter = 0
    temp_chunks_per_section = []

    # Chunk theo từng section với semantic boundaries
    for (stype, snum, scontent) in sections:
        smalls = _chunk_by_semantic_boundaries(scontent, chunk_size, chunk_overlap)
        temp_chunks_per_section.append((stype, snum, smalls))
        total_chunks_counter += len(smalls)

    # Duyệt lần 2 để gán chỉ số chunk toàn cục & metadata
    running_idx = 0
    for (stype, snum, smalls) in temp_chunks_per_section:
        for i, txt in enumerate(smalls):
            out.append({
                "text": txt,
                "chunk_index": running_idx,
                "total_chunks": total_chunks_counter,
                "section_type": stype,
                "section_number": snum,
                "is_complete_section": (len(smalls) == 1),
                "token_count": count_tokens(txt),
                "word_count": len(txt.split())
            })
            running_idx += 1

    return out

# ---------- Embeddings ----------
def get_embeddings(texts: List[str], batch_size: int = 100) -> List[List[float]]:
    if client is None:
        raise Exception("OpenAI client is not initialized")
    all_embeddings = []
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i+batch_size]
        if len(texts) > batch_size:
            st.progress((i + len(batch)) / len(texts), text=f"Generating embeddings... {i + len(batch)}/{len(texts)}")
        resp = client.embeddings.create(model="text-embedding-3-small", input=batch)
        all_embeddings.extend([d.embedding for d in resp.data])
        if i + batch_size < len(texts):
            time.sleep(0.1)
    return all_embeddings
